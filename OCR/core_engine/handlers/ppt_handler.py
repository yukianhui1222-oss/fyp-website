"""
ppt_handler.py
──────────────
PowerPoint (.pptx) 文档处理器，基于 python-pptx。

处理策略：
  - 遍历每张 Slide 的所有 Shape：
      * 文本框 (TEXT_BOX, TITLE) → 提取文字。
      * 图片 (PICTURE) → 提取至内存，送 OCR。
  - 每张幻灯片对应输出一个 page_num。
"""

import io
import logging
import os
import tempfile
from typing import Any, Callable, Dict, List, Optional

from PIL import Image

logger = logging.getLogger(__name__)

# pptx shape 类型枚举（python-pptx MSO_SHAPE_TYPE）
_MSO_PICTURE = 13    # MSO_SHAPE_TYPE.PICTURE
_PP_PLACEHOLDER = 14  # MSO_SHAPE_TYPE.PLACEHOLDER


class PPTHandler:
    """
    PowerPoint 文档处理器。

    Parameters
    ----------
    ocr_processor : OCRProcessor
        已初始化的 OCR 处理器实例（由 DocumentParser 传入）
    """

    def __init__(self, ocr_processor: Any) -> None:
        self.ocr = ocr_processor

    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """Use PowerPoint OLE automation or LibreOffice to convert .ppt to .pptx"""
        # Try win32com first on Windows
        if os.name == 'nt':
            try:
                import win32com.client
                logger.info("Using PowerPoint OLE automation for conversion: %s", ppt_path)
                abs_ppt_path = os.path.abspath(ppt_path).replace('/', '\\')
                temp_dir = tempfile.gettempdir()
                base_name = os.path.basename(ppt_path)
                tmp_pptx_path = os.path.join(temp_dir, f"conv_{base_name}x")
                abs_pptx_path = os.path.abspath(tmp_pptx_path).replace('/', '\\')

                # Ensure CoInitialize
                import pythoncom
                pythoncom.CoInitialize()

                powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
                # Open: FileName, ReadOnly, Untitled, WithWindow
                presentation = powerpoint.Presentations.Open(abs_ppt_path, ReadOnly=True, WithWindow=False)
                presentation.SaveAs(abs_pptx_path, 24) # 24 = ppSaveAsDefault (.pptx)
                presentation.Close()
                powerpoint.Quit()
                return abs_pptx_path
            except Exception as e:
                logger.warning("PowerPoint OLE automation failed or pywin32 missing: %s. Trying LibreOffice...", e)
        
        # Try LibreOffice conversion
        try:
            import subprocess
            import shutil
            soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
            if not soffice_path and os.name == 'nt':
                standard_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
                ]
                for p in standard_paths:
                    if os.path.exists(p):
                        soffice_path = p
                        break
            
            if not soffice_path:
                raise RuntimeError("LibreOffice/soffice executable not found in system PATH.")
                
            temp_dir = tempfile.gettempdir()
            cmd = [soffice_path, "--headless", "--convert-to", "pptx", ppt_path, "--outdir", temp_dir]
            logger.info("Running LibreOffice command: %s", " ".join(cmd))
            
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=30)
            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice failed with code {result.returncode}: {result.stderr}")
                
            base_name = os.path.basename(ppt_path)
            root_name, _ = os.path.splitext(base_name)
            output_filename = f"{root_name}.pptx"
            output_path = os.path.join(temp_dir, output_filename)
            
            if not os.path.exists(output_path):
                raise RuntimeError("Converted file not found at expected path.")
                
            return output_path
        except Exception as e:
            logger.error("All conversion methods failed: %s", e)
            raise RuntimeError(
                "System missing conversion libraries. Please ensure LibreOffice is installed (listed in packages.txt on Streamlit Cloud) or upload .pptx files directly. Error details: " + str(e)
            )

    # ------------------------------------------------------------------
    def process(self, file_path: str, progress_callback: Optional[Callable] = None) -> List[Dict[str, Any]]:
        """
        Parses .pptx files and returns a structured list of results ordered by slide number.

        Parameters
        ----------
        file_path : str
            The absolute path to the PPTX file (.pptx)
        progress_callback : Optional[callable]
            Progress callback
        """
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
        from pptx.util import Emu

        is_legacy = file_path.lower().endswith(".ppt")
        temp_pptx = None

        if is_legacy:
            try:
                temp_pptx = self._convert_ppt_to_pptx(file_path)
                file_path = temp_pptx
            except Exception as e:
                return [
                    {
                        "page_num": 1,
                        "page_type": "error",
                        "text_blocks": [
                            {
                                "text": f"⚠️ Conversion Error: Unable to convert .ppt to .pptx. {e}",
                                "bbox": None,
                                "source": "error",
                                "confidence": 0.0,
                            }
                        ],
                        "table_html": None,
                    }
                ]

        try:
            prs = Presentation(file_path)
        except Exception as e:
            logger.error("Unable to open PPT file: %s", e)
            return [
                {
                    "page_num": 1,
                    "page_type": "error",
                    "text_blocks": [
                        {
                            "text": f"⚠️ Unable to parse this PPT file. The file may be corrupted or in an unsupported format. Error details: {e}",
                            "bbox": None,
                            "source": "error",
                            "confidence": 0.0,
                        }
                    ],
                    "table_html": None,
                }
            ]

        logger.info("PPT has %d slides, starting analysis: %s", len(prs.slides), file_path)

        pages_result: List[Dict[str, Any]] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            native_blocks: List[Dict[str, Any]] = []
            slide_images: List[Image.Image] = []

            for shape in slide.shapes:
                # ── 文本框处理 ──────────────────────────────────────────
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # 尝试获取 shape 的边界框（EMU 转 pt）
                            try:
                                bbox = [
                                    shape.left / 914400 * 72,
                                    shape.top / 914400 * 72,
                                    (shape.left + shape.width) / 914400 * 72,
                                    (shape.top + shape.height) / 914400 * 72,
                                ]
                            except Exception:  # noqa: BLE001
                                bbox = None

                            native_blocks.append(
                                {
                                    "text": text,
                                    "bbox": bbox,
                                    "source": "native",
                                    "confidence": 1.0,
                                }
                            )

                # ── 表格处理 ──────────────────────────────────────────
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                native_blocks.append(
                                    {
                                        "text": cell_text,
                                        "bbox": None,
                                        "source": "native",
                                        "confidence": 1.0,
                                    }
                                )

                # ── 图片处理 ──────────────────────────────────────────
                if shape.shape_type == _MSO_PICTURE:
                    try:
                        img_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                        slide_images.append(img)
                    except Exception as exc:  # noqa: BLE001
                        logger.warning("幻灯片 %d 图片提取失败: %s", slide_num, exc)

            # ── 并发 OCR ──────────────────────────────────────────────
            ocr_blocks: List[Dict[str, Any]] = []
            if slide_images:
                ocr_results = self.ocr.process_images(slide_images, progress_callback=progress_callback)
                for res in ocr_results:
                    if res.get("error"):
                        logger.warning(
                            "Slide %d image OCR failed (index=%d): %s",
                            slide_num, res["index"], res["error"]
                        )
                    ocr_blocks.extend(res["blocks"])

            all_blocks = native_blocks + ocr_blocks
            pages_result.append(
                {
                    "page_num": slide_num,
                    "page_type": "native",
                    "text_blocks": all_blocks,
                    "table_html": None,
                }
            )

        # ── 清理临时文件 ──────────────────────────────────────────────
        if temp_pptx and os.path.exists(temp_pptx):
            try:
                os.remove(temp_pptx)
                logger.info("Cleaned up temporary conversion file: %s", temp_pptx)
            except Exception as e:
                logger.warning("Failed to clean up temporary file: %s", e)

        logger.info("PPT analysis complete, total slides: %d", len(pages_result))
        return pages_result
